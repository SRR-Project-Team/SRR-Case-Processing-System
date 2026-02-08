#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
多类型文件处理器模块

本模块提供各种文件类型的内容提取功能，用于RAG知识库构建。

支持的文件类型：
1. Excel (.xlsx, .xls)
2. Word (.docx)
3. PowerPoint (.pptx)
4. PDF (.pdf)
5. 文本文件 (.txt)
6. CSV (.csv)
7. 图片 (.jpg, .png, .gif等)

作者: Project3 Team
版本: 1.0
"""

import os
from typing import Optional, Dict, Any


def process_excel(file_path: str) -> str:
    """
    处理Excel文件，提取所有sheet的内容
    
    Args:
        file_path: Excel文件路径
        
    Returns:
        str: 提取的文本内容
        
    Raises:
        Exception: 处理失败时抛出异常
    """
    try:
        import pandas as pd
        
        excel_file = pd.ExcelFile(file_path)
        content = ""
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(file_path, sheet_name=sheet_name).fillna("")
            content += f"=== Sheet: {sheet_name} ===\n"
            
            # 按行拼接，每行作为一个案例
            for index, row in df.iterrows():
                row_content = f"Case {index+1}:\n"
                for col in df.columns:
                    row_content += f"{col}: {row[col]}\n"
                content += row_content + "\n"
        
        print(f"✅ Excel文件处理成功: {len(content)} 字符")
        return content
        
    except Exception as e:
        raise Exception(f"Excel文件处理失败: {e}")


def process_word(file_path: str) -> str:
    """
    处理Word文档，提取文本内容
    
    Args:
        file_path: Word文件路径
        
    Returns:
        str: 提取的文本内容
        
    Raises:
        Exception: 处理失败时抛出异常
    """
    try:
        from docx import Document
        
        doc = Document(file_path)
        content = ""
        
        # 提取段落文本
        for para in doc.paragraphs:
            if para.text.strip():
                content += para.text + "\n"
        
        # 提取表格内容
        for table in doc.tables:
            content += "\n=== 表格 ===\n"
            for row in table.rows:
                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                content += row_text + "\n"
        
        print(f"✅ Word文件处理成功: {len(content)} 字符")
        return content
        
    except ImportError:
        raise Exception("python-docx库未安装，请运行: pip install python-docx")
    except Exception as e:
        raise Exception(f"Word文件处理失败: {e}")


def process_powerpoint(file_path: str) -> str:
    """
    处理PowerPoint文件，提取幻灯片文本内容
    
    Args:
        file_path: PowerPoint文件路径
        
    Returns:
        str: 提取的文本内容
        
    Raises:
        Exception: 处理失败时抛出异常
    """
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        content = ""
        
        for slide_num, slide in enumerate(prs.slides, 1):
            content += f"\n=== 幻灯片 {slide_num} ===\n"
            
            # 提取形状中的文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content += shape.text + "\n"
                
                # 提取表格内容
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells])
                        content += row_text + "\n"
        
        print(f"✅ PowerPoint文件处理成功: {len(content)} 字符")
        return content
        
    except ImportError:
        raise Exception("python-pptx库未安装，请运行: pip install python-pptx")
    except Exception as e:
        raise Exception(f"PowerPoint文件处理失败: {e}")


def process_pdf(file_path: str) -> str:
    """
    处理PDF文件，提取文本内容
    
    Args:
        file_path: PDF文件路径
        
    Returns:
        str: 提取的文本内容
        
    Raises:
        Exception: 处理失败时抛出异常
    """
    try:
        import pdfplumber
        
        content = ""
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    content += f"\n=== 页面 {page_num} ===\n"
                    content += page_text + "\n"
        
        print(f"✅ PDF文件处理成功: {len(content)} 字符")
        return content
        
    except Exception as e:
        raise Exception(f"PDF文件处理失败: {e}")


def process_txt(file_path: str) -> str:
    """
    处理文本文件
    
    Args:
        file_path: 文本文件路径
        
    Returns:
        str: 文件内容
        
    Raises:
        Exception: 处理失败时抛出异常
    """
    try:
        # 尝试多种编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'utf-16', 'latin1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                print(f"✅ 文本文件处理成功 ({encoding}): {len(content)} 字符")
                return content
            except UnicodeDecodeError:
                continue
        
        # 如果所有编码都失败，使用错误忽略模式
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        print(f"⚠️ 文本文件处理成功 (错误忽略模式): {len(content)} 字符")
        return content
        
    except Exception as e:
        raise Exception(f"文本文件处理失败: {e}")


def process_csv(file_path: str) -> str:
    """
    处理CSV文件
    
    Args:
        file_path: CSV文件路径
        
    Returns:
        str: 提取的文本内容
        
    Raises:
        Exception: 处理失败时抛出异常
    """
    try:
        import pandas as pd
        
        # 尝试多种编码读取CSV
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin1']
        df = None
        
        for encoding in encodings:
            try:
                df = pd.read_csv(file_path, encoding=encoding)
                break
            except UnicodeDecodeError:
                continue
        
        if df is None:
            raise Exception("无法使用任何编码读取CSV文件")
        
        # 填充空值
        df = df.fillna("")
        
        # 转换为文本格式
        content = "=== CSV数据 ===\n"
        for index, row in df.iterrows():
            row_content = f"行 {index+1}:\n"
            for col in df.columns:
                row_content += f"{col}: {row[col]}\n"
            content += row_content + "\n"
        
        print(f"✅ CSV文件处理成功: {len(content)} 字符")
        return content
        
    except Exception as e:
        raise Exception(f"CSV文件处理失败: {e}")


def process_image(file_path: str) -> str:
    """
    处理图片文件，使用OCR提取文本
    
    Args:
        file_path: 图片文件路径
        
    Returns:
        str: OCR提取的文本内容
        
    Raises:
        Exception: 处理失败时抛出异常
    """
    try:
        from PIL import Image
        import pytesseract
        
        # 打开图片
        image = Image.open(file_path)
        
        # 使用OCR提取文本（支持中英文）
        text = pytesseract.image_to_string(image, lang='chi_sim+eng')
        
        if not text.strip():
            print("⚠️ 图片中未检测到文本内容")
            return "[图片文件，未检测到文本内容]"
        
        print(f"✅ 图片文件处理成功: {len(text)} 字符")
        return f"=== 图片OCR结果 ===\n{text}"
        
    except ImportError:
        # 如果OCR库未安装，返回占位符
        print("⚠️ pytesseract未安装，无法进行OCR")
        return "[图片文件，OCR功能未启用]"
    except Exception as e:
        print(f"⚠️ 图片处理失败: {e}")
        return f"[图片文件，处理失败: {str(e)}]"


def process_file(file_path: str, file_type: str) -> str:
    """
    根据文件类型自动选择处理器
    
    Args:
        file_path: 文件路径
        file_type: 文件类型（excel, word, powerpoint, pdf, txt, csv, image）
        
    Returns:
        str: 提取的文本内容
        
    Raises:
        Exception: 处理失败时抛出异常
    """
    processors = {
        'excel': process_excel,
        'word': process_word,
        'powerpoint': process_powerpoint,
        'pdf': process_pdf,
        'txt': process_txt,
        'csv': process_csv,
        'image': process_image
    }
    
    processor = processors.get(file_type.lower())
    
    if processor is None:
        raise Exception(f"不支持的文件类型: {file_type}")
    
    return processor(file_path)


def get_file_metadata(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    获取文件元数据
    
    Args:
        file_path: 文件路径
        file_type: 文件类型
        
    Returns:
        Dict[str, Any]: 文件元数据字典
    """
    metadata = {
        'file_type': file_type,
        'file_size': os.path.getsize(file_path)
    }
    
    try:
        if file_type == 'excel':
            import pandas as pd
            excel_file = pd.ExcelFile(file_path)
            metadata['sheet_count'] = len(excel_file.sheet_names)
            metadata['sheet_names'] = excel_file.sheet_names
            
        elif file_type == 'word':
            from docx import Document
            doc = Document(file_path)
            metadata['paragraph_count'] = len(doc.paragraphs)
            metadata['table_count'] = len(doc.tables)
            
        elif file_type == 'powerpoint':
            from pptx import Presentation
            prs = Presentation(file_path)
            metadata['slide_count'] = len(prs.slides)
            
        elif file_type == 'pdf':
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                metadata['page_count'] = len(pdf.pages)
                
        elif file_type == 'image':
            from PIL import Image
            img = Image.open(file_path)
            metadata['width'] = img.width
            metadata['height'] = img.height
            metadata['format'] = img.format
            
    except Exception as e:
        print(f"⚠️ 获取元数据时出错: {e}")
    
    return metadata


def detect_file_type_from_extension(filename: str) -> str:
    """
    根据文件扩展名检测文件类型
    
    Args:
        filename: 文件名
        
    Returns:
        str: 文件类型（excel, word, powerpoint, pdf, txt, csv, image, unknown）
    """
    ext = os.path.splitext(filename)[1].lower()
    
    type_mapping = {
        '.xlsx': 'excel',
        '.xls': 'excel',
        '.docx': 'word',
        '.doc': 'word',
        '.pptx': 'powerpoint',
        '.ppt': 'powerpoint',
        '.pdf': 'pdf',
        '.txt': 'txt',
        '.csv': 'csv',
        '.jpg': 'image',
        '.jpeg': 'image',
        '.png': 'image',
        '.gif': 'image',
        '.bmp': 'image',
        '.tiff': 'image'
    }
    
    return type_mapping.get(ext, 'unknown')
